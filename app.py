"""Classroom Report — single file: config, parsers, analytics, LLM, reports, LangGraph, FastAPI, CLI, Streamlit."""
from __future__ import annotations

import argparse
import base64
import json
import os
import sys
import tempfile
import urllib.error
import urllib.request
from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any, List, NotRequired, Optional, Tuple, TypedDict, Union

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from langgraph.graph import END, StateGraph
from pptx import Presentation
from pydantic import BaseModel, Field

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None  # type: ignore

# === config ===
import os

# Ollama
OLLAMA_HOST = os.environ.get("OLLAMA_HOST", "http://localhost:11434")
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "llama3.2")

# Tier thresholds (percentiles): top %, middle %, low %
TIER_TOP_PCT = 20
TIER_AVERAGE_PCT = 60
TIER_LOW_PCT = 20

# File size limits (bytes)
MAX_SLIDES_SIZE_BYTES = 50 * 1024 * 1024   # 50 MB
MAX_EXCEL_SIZE_BYTES = 10 * 1024 * 1024   # 10 MB

# Allowed extensions
ALLOWED_SLIDE_EXTENSIONS = (".pptx", ".pdf")
ALLOWED_EXCEL_EXTENSIONS = (".xlsx", ".xls")

# Excel column name for student identity (case-insensitive match)
STUDENT_NAME_COLUMN = "Student Name"

# Neutral tier labels for homework doc (pedagogical guardrail)
TIER_LABELS = {
    "top": "Extension",
    "average": "Core",
    "low": "Support",
}

# === parsers/responses ===
from io import BytesIO
"""Load and normalize poll responses from Excel; compute scores and tiers."""

import pandas as pd



def _find_student_name_column(df: pd.DataFrame) -> str:
    """Find column that matches STUDENT_NAME_COLUMN (case-insensitive)."""
    for c in df.columns:
        if str(c).strip().lower() == STUDENT_NAME_COLUMN.lower():
            return c
    raise ValueError(
        f"Excel must have a column named '{STUDENT_NAME_COLUMN}'. "
        "Expected format: Student Name + Q1, Q2, ... or Student Name + Q1_Selected/Q1_Correct pairs."
    )


def _detect_selected_correct_pairs(df: pd.DataFrame, name_col: str) -> Optional[List[Tuple[str, str, str]]]:
    """
    If Excel has Qn_Selected and Qn_Correct columns, return list of (question_label, selected_col, correct_col).
    E.g. [("Q1", "Q1_Selected", "Q1_Correct"), ("Q2", "Q2_Selected", "Q2_Correct"), ...]
    Uses actual df column names for lookup.
    """
    col_map = {str(c).strip(): c for c in df.columns}
    cols = list(col_map.keys())
    selected = [c for c in cols if c.endswith("_Selected") or c.endswith("_selected")]
    if not selected:
        return None
    pairs = []
    for c in selected:
        base = c.replace("_Selected", "").replace("_selected", "")
        correct_name = base + "_Correct"
        correct_alt = base + "_correct"
        if correct_name in cols:
            pairs.append((base, col_map[c], col_map[correct_name]))
        elif correct_alt in cols:
            pairs.append((base, col_map[c], col_map[correct_alt]))
    if not pairs:
        return None
    return sorted(pairs, key=lambda x: x[0])


def _get_question_columns(df: pd.DataFrame, name_col: str) -> List[str]:
    """Return list of question columns (exclude name column and any non-question columns)."""
    exclude = {name_col}
    # Common question patterns: Q1, Q2, Question 1, etc.
    q_cols = []
    for c in df.columns:
        if c in exclude:
            continue
        if isinstance(c, str) and (c.strip().upper().startswith("Q") or "question" in c.strip().lower()):
            if "_Selected" not in c and "_Correct" not in c:
                q_cols.append(c)
        elif isinstance(c, (int, float)) or (isinstance(c, str) and c.strip().replace(".", "").isdigit()):
            q_cols.append(c)
    if not q_cols:
        # Fallback: all columns except name and Email ID / Selected / Correct
        q_cols = [
            c for c in df.columns
            if c not in exclude
            and "_Selected" not in str(c)
            and "_Correct" not in str(c)
        ]
    return q_cols


def load_responses(
    file_or_bytes: BytesIO | bytes | str,
    answer_key: Optional[str] = None,
) -> pd.DataFrame:
    """
    Load Excel and return a normalized DataFrame: Student Name + Q1, Q2, ... with 1=correct, 0=incorrect.

    Supports two formats:
    1. Selected/Correct: columns like Q1_Selected, Q1_Correct (letters A/B/C/D). Score 1 where Selected == Correct.
    2. Wide: Student Name + Q1, Q2, ... with values 1/0 or A/B/C and optional answer_key (or first row as key).
    """
    df = pd.read_excel(file_or_bytes, engine="openpyxl")
    if df.empty or len(df) < 1:
        raise ValueError("Excel file is empty or has no data rows.")
    name_col = _find_student_name_column(df)

    # Format 1: Qn_Selected and Qn_Correct columns (e.g. Q1_Selected, Q1_Correct)
    pairs = _detect_selected_correct_pairs(df, name_col)
    if pairs:
        out = df[[name_col]].copy()
        out.columns = [STUDENT_NAME_COLUMN]
        for q_label, selected_col, correct_col in pairs:
            selected = df[selected_col].astype(str).str.strip().str.upper()
            correct = df[correct_col].astype(str).str.strip().str.upper()
            out[q_label] = (selected == correct).astype(int)
        return out

    # Format 2: wide format (Q1, Q2, ... with 1/0 or answer key)
    q_cols = _get_question_columns(df, name_col)
    if not q_cols:
        raise ValueError(
            "No question columns found. Use Q1, Q2, ... or Q1_Selected/Q1_Correct, Q2_Selected/Q2_Correct, etc."
        )

    # Optional: first row as answer key (if first cell is "key"/"answer" or user provided key)
    key_row = None
    if answer_key:
        key_parts = [x.strip() for x in answer_key.split(",")]
        if len(key_parts) >= len(q_cols):
            key_row = key_parts[: len(q_cols)]
        else:
            key_row = key_parts + [None] * (len(q_cols) - len(key_parts))
    else:
        first_cell = str(df.iloc[0].get(name_col, "")).strip().lower()
        if first_cell in ("key", "answer", "answers", ""):
            first = df.iloc[0]
            key_row = [first.get(c) for c in q_cols]
            df = df.iloc[1:].reset_index(drop=True)

    # Normalize to 0/1 (ensure numeric: Excel often gives "1.0"/"0.0" as strings)
    out = df[[name_col]].copy()
    out.columns = [STUDENT_NAME_COLUMN]
    for j, q in enumerate(q_cols):
        col = df[q].astype(str).str.strip().str.lower()
        if key_row is not None and j < len(key_row):
            correct_val = str(key_row[j]).strip().lower() if key_row[j] is not None else "1"
            out[q] = (col == correct_val).astype(int)
        else:
            # Assume 1 = correct, 0 = incorrect; map common values then coerce to numeric
            mapped = col.replace(
                {"1": 1, "1.0": 1, "yes": 1, "y": 1, "correct": 1, "true": 1, "0": 0, "0.0": 0, "no": 0, "n": 0, "false": 0, "": 0}
            )
            out[q] = pd.to_numeric(mapped, errors="coerce").fillna(0).clip(0, 1).astype(int)
    return out


def normalize_responses(
    df: pd.DataFrame,
) -> Tuple[pd.DataFrame, List[str]]:
    """
    Ensure DataFrame has Student Name and question columns; return (df, question_columns).
    """
    name_col = _find_student_name_column(df)
    q_cols = _get_question_columns(df, name_col)
    return df[[name_col] + q_cols], q_cols

# === parsers/slides_pptx ===
"""Extract text from PowerPoint (.pptx) files."""
import tempfile
from io import BytesIO
from pathlib import Path
from typing import List, Tuple, Union

from pptx import Presentation


def extract_text_from_pptx(path: Union[str, Path, BytesIO]) -> Tuple[str, List[Tuple[int, str]]]:
    """
    Extract all text from a .pptx file (path or BytesIO).
    Returns (full_lecture_text, list of (slide_index_1based, slide_text)).
    """
    if isinstance(path, BytesIO):
        path = BytesIO(path.getvalue())
        path.name = "upload.pptx"
    p = Path(path) if not isinstance(path, BytesIO) else None
    if p and p.suffix.lower() != ".pptx":
        raise ValueError("File must be .pptx")
    if isinstance(path, BytesIO):
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(path.getvalue())
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
        finally:
            Path(tmp_path).unlink(missing_ok=True)
    else:
        prs = Presentation(str(path))
    slides_text: List[Tuple[int, str]] = []
    all_parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        slide_text = "\n".join(parts)
        slides_text.append((i, slide_text))
        if slide_text:
            all_parts.append(slide_text)
    full_text = "\n\n".join(all_parts)
    return full_text, slides_text


def get_poll_slides(slides_text: List[Tuple[int, str]]) -> List[Tuple[int, str]]:
    """Return slides that look like poll/question slides (by title or first line)."""
    poll_keywords = ("poll", "question", "quiz")
    result = []
    for idx, text in slides_text:
        first_line = (text.split("\n")[0] or "").strip().lower()
        if any(kw in first_line for kw in poll_keywords):
            result.append((idx, text))
    return result

# === parsers/slides_pdf ===
"""Extract text from PDF files."""
import tempfile
from io import BytesIO
from pathlib import Path
from typing import List, Tuple, Union

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None  # type: ignore


def extract_text_from_pdf(path: Union[str, Path, BytesIO]) -> Tuple[str, List[Tuple[int, str]]]:
    """
    Extract all text from a PDF file (path or BytesIO). One page = one "slide".
    Returns (full_lecture_text, list of (page_number_1based, page_text)).
    """
    if fitz is None:
        raise ImportError("PyMuPDF is required for PDF support. Install with: pip install PyMuPDF")
    if isinstance(path, BytesIO):
        doc = fitz.open(stream=path.getvalue(), filetype="pdf")
    else:
        path = Path(path)
        if path.suffix.lower() != ".pdf":
            raise ValueError("File must be .pdf")
        doc = fitz.open(str(path))
    slides_text: List[Tuple[int, str]] = []
    all_parts: List[str] = []
    try:
        for i in range(len(doc)):
            page = doc[i]
            text = page.get_text().strip()
            slides_text.append((i + 1, text))
            if text:
                all_parts.append(text)
    finally:
        doc.close()
    full_text = "\n\n".join(all_parts)
    return full_text, slides_text


def get_poll_slides(slides_text: List[Tuple[int, str]]) -> List[Tuple[int, str]]:
    """Return slides/pages that look like poll/question (first line contains keyword)."""
    poll_keywords = ("poll", "question", "quiz")
    result = []
    for idx, text in slides_text:
        first_line = (text.split("\n")[0] or "").strip().lower()
        if any(kw in first_line for kw in poll_keywords):
            result.append((idx, text))
    return result

# === analytics/scoring ===
"""Deterministic scoring and tier assignment from poll response DataFrame."""

import pandas as pd



def compute_scores(
    df: pd.DataFrame,
    question_columns: Optional[List[str]] = None,
) -> pd.DataFrame:
    """
    Compute per-student score = (number correct) / (number answered).
    Blank/missing excluded from denominator. Returns df with 'score' and 'answered' columns.
    """
    name_col = _find_student_name_column(df)
    if question_columns is None:
        question_columns = _get_question_columns(df, name_col)
    q = df[question_columns]
    # Coerce to numeric 0/1 in case Excel left strings (e.g. "1.0")
    q = q.apply(pd.to_numeric, errors="coerce").fillna(0).clip(0, 1)
    answered = q.notna() & (q != "")
    correct = (q == 1).fillna(False)
    n_answered = answered.sum(axis=1).astype(int)
    n_correct = (correct & answered).sum(axis=1).astype(int)
    score = n_correct / n_answered.replace(0, float("nan"))
    out = df[[name_col]].copy()
    out["answered"] = n_answered
    out["score"] = score
    out["score_pct"] = (score * 100).round(1)
    return out


def assign_tiers(
    scores_df: pd.DataFrame,
    top_pct: int = TIER_TOP_PCT,
    average_pct: int = TIER_AVERAGE_PCT,
    low_pct: int = TIER_LOW_PCT,
) -> pd.DataFrame:
    """
    Assign tier by percentile: top_pct = 'top', next average_pct = 'average', rest = 'low'.
    scores_df must have 'score' column. Adds 'tier' and 'rank' columns.
    """
    df = scores_df.copy()
    df = df.dropna(subset=["score"])
    df = df.sort_values("score", ascending=False).reset_index(drop=True)
    n = len(df)
    df["rank"] = range(1, n + 1)
    if n == 0:
        df["tier"] = []
        return df
    top_n = max(1, round(n * top_pct / 100))
    low_n = max(0, round(n * low_pct / 100))
    mid_n = n - top_n - low_n
    tier = ["top"] * top_n + ["average"] * mid_n + ["low"] * low_n
    df["tier"] = tier[:n]
    return df


def get_top_n(ranked_df: pd.DataFrame, n: int = 5) -> pd.DataFrame:
    """Return top N performers (by rank)."""
    name_col = _find_student_name_column(ranked_df)
    if "rank" not in ranked_df.columns:
        return ranked_df.head(n)
    return ranked_df.nsmallest(n, "rank")[[name_col, "score", "score_pct", "rank"]]

# === analytics/charts ===
"""Plotly charts for top 5 performers and per-question engagement."""

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go



def chart_top5(
    top5_df: pd.DataFrame,
    anonymize: bool = False,
    name_col: Optional[str] = None,
) -> go.Figure:
    """Bar chart of top 5 performers: name (or Student 1..5) vs score %."""
    if top5_df.empty:
        fig = go.Figure()
        fig.add_annotation(text="No data", xref="paper", yref="paper", x=0.5, y=0.5, showarrow=False)
        return fig
    name_col = name_col or STUDENT_NAME_COLUMN
    df = top5_df.copy()
    if anonymize:
        df["Display"] = [f"Student {i+1}" for i in range(len(df))]
    else:
        df["Display"] = df[name_col].astype(str)
    fig = px.bar(
        df,
        x="Display",
        y="score_pct",
        title="Top 5 performers (poll score %)",
        labels={"score_pct": "Score (%)", "Display": "Student"},
        text_auto=".1f",
    )
    fig.update_layout(xaxis_tickangle=-45, showlegend=False)
    return fig


def chart_engagement(
    responses_df: pd.DataFrame,
    question_columns: List[str],
) -> go.Figure:
    """Bar chart: per-question engagement (count answered) and optionally % correct."""
    q = responses_df[question_columns].apply(pd.to_numeric, errors="coerce").fillna(0)
    answered = q.notna() & (q != "")
    count_answered = answered.sum().astype(int)
    correct = (q == 1) & answered
    pct_correct = correct.sum() / count_answered.replace(0, float("nan")) * 100
    df = pd.DataFrame({
        "Question": question_columns,
        "Answered": count_answered.values,
        "Correct %": pct_correct.values.round(1),
    })
    fig = go.Figure()
    fig.add_trace(go.Bar(name="Responses", x=df["Question"], y=df["Answered"], text=df["Answered"].astype(int), textposition="outside"))
    fig.add_trace(go.Scatter(x=df["Question"], y=df["Correct %"], name="Correct %", mode="lines+markers", yaxis="y2"))
    fig.update_layout(
        title="Poll engagement by question",
        xaxis_title="Question",
        yaxis=dict(title="Number of responses"),
        yaxis2=dict(title="Correct %", overlaying="y", side="right", range=[0, 105]),
        barmode="group",
        showlegend=True,
    )
    return fig

# === pipeline/ingest ===
"""Parse slides and Excel; compute analytics — used by Streamlit, API, and LangGraph."""

from dataclasses import dataclass
from io import BytesIO
from typing import Optional

import pandas as pd



@dataclass
class AnalyticsBundle:
    """DataFrames from a completed analytics pass."""

    responses_df: pd.DataFrame
    question_columns: list[str]
    scores_df: pd.DataFrame
    ranked_df: pd.DataFrame
    top5_df: pd.DataFrame
    tier_counts: dict[str, int]


def parse_slides_bytes(filename: str, data: bytes) -> tuple[str, str]:
    """
    Return (lecture_text, poll_questions_text) from slide bytes.
    filename must include extension (.pptx or .pdf).
    """
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    buf = BytesIO(data)
    buf.name = filename
    if ext == "pptx":
        # use extract_text_from_pptx, get_poll_slides below

        full_text, slides_text = extract_text_from_pptx(buf)
    elif ext == "pdf":
        

        full_text, slides_text = extract_text_from_pdf(buf)
    else:
        return "", ""

    poll_slides = get_poll_slides(slides_text)
    poll_text = "\n\n".join([f"Slide {i}: {t}" for i, t in poll_slides]) if poll_slides else ""
    return full_text, poll_text


def parse_responses_bytes(
    excel_bytes: bytes,
    filename: str = "responses.xlsx",
    answer_key: Optional[str] = None,
) -> tuple[pd.DataFrame, list[str]]:
    """Load Excel bytes and return (df, question_columns). Raises on error."""
    buf = BytesIO(excel_bytes)
    buf.name = filename
    df = load_responses(buf, answer_key=answer_key)
    df, q_cols = normalize_responses(df)
    return df, q_cols


def run_analytics(df: pd.DataFrame, question_columns: list[str]) -> AnalyticsBundle:
    """Compute scores, tiers, top 5, and tier counts."""
    scores_df = compute_scores(df, question_columns)
    ranked_df = assign_tiers(scores_df)
    top5_df = get_top_n(ranked_df, 5)
    tier_counts = ranked_df["tier"].value_counts().to_dict() if "tier" in ranked_df.columns else {}
    return AnalyticsBundle(
        responses_df=df,
        question_columns=question_columns,
        scores_df=scores_df,
        ranked_df=ranked_df,
        top5_df=top5_df,
        tier_counts=tier_counts,
    )


def assert_slides_size(num_bytes: int) -> None:
    if num_bytes > MAX_SLIDES_SIZE_BYTES:
        raise ValueError(f"Slides file too large. Max {MAX_SLIDES_SIZE_BYTES // (1024 * 1024)} MB.")


def assert_excel_size(num_bytes: int) -> None:
    if num_bytes > MAX_EXCEL_SIZE_BYTES:
        raise ValueError(f"Excel file too large. Max {MAX_EXCEL_SIZE_BYTES // (1024 * 1024)} MB.")

# === llm/ollama_client ===
"""Ollama client: prompt to string with connection checks."""
import os
from typing import Optional



def get_ollama_host() -> str:
    return os.environ.get("OLLAMA_HOST", OLLAMA_HOST)


def check_ollama_available(host: Optional[str] = None) -> tuple[bool, str]:
    """
    Ping Ollama (e.g. GET /api/tags). Returns (success, message).
    """
    import urllib.request
    import urllib.error

    base = (host or get_ollama_host()).rstrip("/")
    url = f"{base}/api/tags"
    try:
        req = urllib.request.Request(url, method="GET")
        with urllib.request.urlopen(req, timeout=5) as r:
            if r.status == 200:
                return True, "Ollama is available."
            return False, f"Ollama returned status {r.status}"
    except urllib.error.URLError as e:
        return False, f"Cannot reach Ollama: {e.reason}. Start Ollama and pull a model (e.g. ollama pull llama3.2)."
    except Exception as e:
        return False, str(e)


def prompt_ollama(
    prompt: str,
    model: str = "llama3.2",
    host: Optional[str] = None,
    system: Optional[str] = None,
) -> str:
    """
    Send prompt to Ollama chat API; return assistant message content.
    Raises RuntimeError if Ollama is unavailable or returns an error.
    """
    try:
        from ollama import chat
    except ImportError:
        raise ImportError("ollama package required. Install with: pip install ollama")

    base = host or get_ollama_host()
    # ollama package uses OLLAMA_HOST env
    env_host = os.environ.get("OLLAMA_HOST")
    try:
        if base:
            os.environ["OLLAMA_HOST"] = base
        messages = []
        if system:
            messages.append({"role": "system", "content": system})
        messages.append({"role": "user", "content": prompt})
        response = chat(model=model, messages=messages)
        content = getattr(response, "message", None) or (response if isinstance(response, dict) else {})
        if hasattr(content, "content"):
            return content.content or ""
        if isinstance(content, dict) and "content" in content:
            return content.get("content", "") or ""
        return str(content)
    finally:
        if env_host is not None:
            os.environ["OLLAMA_HOST"] = env_host
        elif "OLLAMA_HOST" in os.environ and not base:
            del os.environ["OLLAMA_HOST"]


class OllamaClient:
    """Thin wrapper for summary and homework generation."""

    def __init__(self, model: str = "llama3.2", host: Optional[str] = None):
        self.model = model
        self.host = host or get_ollama_host()

    def available(self) -> tuple[bool, str]:
        return check_ollama_available(self.host)

    def generate_topic_summary(
        self,
        lecture_text: str,
        poll_questions_text: str = "",
        extra_system: Optional[str] = None,
    ) -> str:
        """Generate 1–2 paragraph topic summary from lecture and optional poll questions."""
        system = "You are a helpful assistant that writes concise, factual summaries for teachers."
        if extra_system and extra_system.strip():
            system = system + "\n\n" + extra_system.strip()
        prompt = (
            "Based on the following lecture content, write a short topic summary (1–2 paragraphs) "
            "suitable for a class report. Be concise and focus on the main concepts discussed.\n\n"
            "Lecture content:\n" + lecture_text
        )
        if poll_questions_text:
            prompt += "\n\nPoll questions covered:\n" + poll_questions_text
        prompt += "\n\nWrite only the summary, no headings."
        return prompt_ollama(prompt, model=self.model, host=self.host, system=system)

    def generate_differentiated_homework(
        self,
        topic_summary: str,
        tier_counts: dict[str, int],
        question_specs: Optional[list[dict]] = None,
        levels: Optional[list[str]] = None,
        extra_system: Optional[str] = None,
    ) -> str:
        """
        Generate differentiated homework based on lecture/PPT content.
        tier_counts: {'top': n, 'average': n, 'low': n}.
        question_specs: optional list of {"type": "MCQ"|"Fill in the blanks"|"Subjective questions", "count": n}.
        levels: list of level names to generate, e.g. ["Extension", "Core", "Support"] or ["Core"] only.
        Answer key for MCQs must be at the end, not inline.
        """
        if not (topic_summary or "").strip():
            return (
                "Extension\n\nBased on today's class, complete 2–3 extension tasks that go beyond the lesson.\n\n"
                "Core\n\nComplete the standard practice set based on today's topic.\n\n"
                "Support\n\nWork through the guided practice and review the key points from class.\n\n"
                "Answer key\n\n(MCQ answers listed here when generated.)"
            )
        specs = question_specs or [
            {"type": "MCQ", "count": 2},
            {"type": "Fill in the blanks", "count": 2},
            {"type": "Subjective questions", "count": 1},
        ]
        levels = levels or ["Extension", "Core", "Support"]
        spec_text = ", ".join(f"{s['count']} {s['type']}" for s in specs)
        levels_text = ", ".join(levels)

        system = (
            "You are a helpful assistant for teachers. You generate concrete homework activities. "
            "Do not mention student names or performance levels. "
            "Use only the section headings requested (Extension, Core, Support). "
            "For MCQs: give question and 4 options (A, B, C, D) only — do NOT write the correct answer next to the question. "
            "Put all MCQ correct answers in a separate 'Answer key' section at the very end of your output."
        )
        if extra_system and extra_system.strip():
            system = system + "\n\n" + extra_system.strip()
        prompt = (
            "Use the following lecture/topic content to generate differentiated homework.\n\n"
            "---\nLECTURE / TOPIC CONTENT\n---\n\n"
            + topic_summary.strip()
            + "\n\n---\n"
            "Generate ONLY these sections (in this order): " + levels_text + ".\n\n"
            f"For EACH of these sections, generate exactly: {spec_text}.\n\n"
            "Requirements:\n"
            "- MCQ: write the question and four options (A, B, C, D). Do NOT indicate the correct answer in the question. "
            "At the very end of your output, add a section titled exactly 'Answer key' (or 'Answer Key') and list only the correct answers for every MCQ, "
            "e.g. 'Extension: 1. A, 2. C | Core: 1. B, 2. D | Support: 1. A' or similar so the teacher can use it for grading.\n"
            "- Fill in the blanks: write a sentence with ____ for the blank; base the word on the lecture.\n"
            "- Subjective questions: open-ended questions for a short paragraph or list answer.\n"
            "Extension: slightly harder or extension-oriented. Core: standard difficulty. Support: simpler and more guided.\n\n"
            f"Output format: put each of these headings on its own line: {levels_text}. Under each heading list the questions. "
            "After all level sections, add the 'Answer key' section. Do not output anything before the first heading or after Answer key."
        )
        return prompt_ollama(prompt, model=self.model, host=self.host, system=system)

# === reports/summary_doc ===
"""Build summary report .docx from Ollama output."""



def build_summary_docx(
    summary_text: str,
    title: str = "Class Topic Summary",
    date_placeholder: bool = True,
) -> bytes:
    """
    Create a Word document with title, optional date, and summary section.
    Returns document as bytes for download.
    """
    doc = Document()
    # Title
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if date_placeholder:
        doc.add_paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d')}")
    doc.add_paragraph()
    # Summary
    doc.add_heading("Summary", level=1)
    for para in summary_text.strip().split("\n\n"):
        if para.strip():
            p = doc.add_paragraph(para.strip())
            p.paragraph_format.space_after = Pt(6)
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()

# === reports/homework_doc ===
"""Build differentiated homework .docx from Ollama tiered output."""



def build_homework_docx(
    homework_text: str,
    title: str = "Differentiated Homework",
    date_placeholder: bool = True,
) -> bytes:
    """
    Create a Word document with title, date, and three sections (Extension / Core / Support).
    homework_text should be the raw LLM output with headings Extension, Core, Support.
    Returns document as bytes.
    """
    doc = Document()
    doc.add_heading(title, level=0)
    if date_placeholder:
        doc.add_paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d')}")
    doc.add_paragraph()

    raw = (homework_text or "").strip()
    if not raw:
        doc.add_paragraph(
            "No homework content was generated. Try generating the topic summary first, then generate homework again."
        )
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.read()

    # Parse into (heading, content) and render level sections first, Answer key last
    lines = raw.split("\n")
    level_headers = ("extension", "core", "support")
    answer_key_headers = ("answer key", "answerkey")
    current_heading = None
    current_lines = []
    answer_key_content = []  # collected separately, rendered at the end

    def flush(heading: str, content_lines: list):
        if not heading:
            return
        block = "\n".join(content_lines).strip()
        if not block:
            return
        doc.add_heading(heading, level=1)
        for p in block.split("\n\n"):
            if p.strip():
                doc.add_paragraph(p.strip())

    for line in lines:
        stripped = line.strip()
        lower = stripped.lstrip("#* ").lower()
        is_answer_key = any(lower.startswith(ak) for ak in answer_key_headers) or lower == "answer key"
        is_level = any(lower.startswith(h) for h in level_headers)

        if is_answer_key:
            # Flush current level section, then collect Answer key content (rendered later)
            if current_heading:
                flush(current_heading, current_lines)
                current_lines.clear()
            current_heading = "Answer key"
            rest = stripped.lstrip("#* ")
            for ak in ("Answer key", "Answer Key", "answer key"):
                if rest.lower().startswith(ak.lower()):
                    rest = rest[len(ak) :].lstrip(":.- ")
                    break
            if rest:
                answer_key_content.append(rest)
        elif is_level:
            if current_heading and current_heading != "Answer key":
                flush(current_heading, current_lines)
                current_lines.clear()
            elif current_heading == "Answer key":
                # Already in answer key; keep appending to answer_key_content
                pass
            current_heading = (
                "Extension" if lower.startswith("extension") else "Core" if lower.startswith("core") else "Support"
            )
            if current_heading != "Answer key":
                rest = stripped.lstrip("#* ")
                for h in ("Extension", "Core", "Support"):
                    if rest.lower().startswith(h.lower()):
                        rest = rest[len(h) :].lstrip(":.- ")
                        break
                if rest:
                    current_lines.append(rest)
        else:
            if current_heading == "Answer key":
                answer_key_content.append(line)
            elif current_heading:
                current_lines.append(line)

    if current_heading and current_heading != "Answer key":
        flush(current_heading, current_lines)

    # Answer key at the end
    if answer_key_content:
        flush("Answer key", answer_key_content)

    # Fallback: no structured sections parsed
    if len(doc.paragraphs) <= 2:
        doc.add_heading("Homework", level=1)
        for para in raw.split("\n\n"):
            if para.strip():
                doc.add_paragraph(para.strip())

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()

# === agent/state ===
"""LangGraph state: dynamic optional fields via total=False."""

from typing import NotRequired, TypedDict


class ClassroomState(TypedDict, total=False):
    """Workflow state; fields appear as the run progresses."""

    # Loaded from agent.md / skills.md
    agent_context: str
    skills_context: str

    # Inputs
    slides_bytes: bytes
    excel_bytes: bytes
    slides_filename: str
    excel_filename: str
    answer_key: str | None
    ollama_model: str
    want_summary: NotRequired[bool]
    want_homework: NotRequired[bool]
    anonymize: NotRequired[bool]
    homework_levels: NotRequired[list[str]]
    question_specs: NotRequired[list[dict]]

    # Ingest + analytics
    ingest_ok: bool
    lecture_text: str
    poll_questions_text: str
    tier_counts: dict[str, int]
    ranked_preview: list[dict]
    question_columns: list[str]
    charts: dict[str, str]

    # LLM outputs
    summary_text: str
    homework_text: str
    summary_docx_bytes: bytes
    homework_docx_bytes: bytes

    errors: list[str]

# === agent/loaders ===
"""Load agent.md (memory/persona) and skills.md (workflow instructions)."""

from pathlib import Path


def project_root() -> Path:
    """Repository root (directory containing this file)."""
    return Path(__file__).resolve().parent


def default_agent_md() -> str:
    return """# Agent memory

You support teachers with **local-only** classroom analytics and report generation.

## Principles

- Prefer concise, factual outputs suitable for class reports.
- Respect file formats: slides as PPT/PDF, poll data as Excel with `Student Name`.
- Do not send data to the cloud; Ollama runs locally.

## Disclaimer

Analytics and tiers are based only on poll responses; use professional judgment when assigning homework.
"""


def default_skills_md() -> str:
    return """# Skills

## Ingest

1. Accept lecture slides (`.pptx` or `.pdf`) and poll responses (`.xlsx`).
2. Parse slides for full text and poll-tagged slides (titles containing Poll/Question/Quiz).
3. Parse Excel: wide format (`Q1`, `Q2`, …) or Selected/Correct pairs (`Q1_Selected`, `Q1_Correct`).

## Analytics

1. Score students, assign tiers (Extension / Core / Support), compute top performers.
2. Build engagement charts from question columns.

## Reports

1. **Topic summary** — short class summary from slide content (optional poll context).
2. **Homework** — differentiated Extension / Core / Support; no student names in LLM prompts for homework (only topic + tier counts).

## Guardrails

- Homework prompts: include topic text and tier counts only — not student names.
- MCQ answers belong in a final Answer key section, not inline.
"""


def _read_or_create(path: Path, default: str) -> str:
    if path.is_file():
        return path.read_text(encoding="utf-8")
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(default, encoding="utf-8")
    return default


def load_agent_md(root: Path | None = None) -> str:
    p = (root or project_root()) / "agent.md"
    return _read_or_create(p, default_agent_md())


def load_skills_md(root: Path | None = None) -> str:
    p = (root or project_root()) / "skills.md"
    return _read_or_create(p, default_skills_md())


def combine_agent_skills(agent_context: str, skills_context: str) -> str:
    parts: list[str] = []
    if (agent_context or "").strip():
        parts.append(agent_context.strip())
    if (skills_context or "").strip():
        parts.append(skills_context.strip())
    return "\n\n---\n\n".join(parts) if parts else ""

# === agent/graph ===


def _node_load_context(state: ClassroomState) -> dict[str, Any]:
    agent = load_agent_md()
    skills = load_skills_md()
    return {"agent_context": agent, "skills_context": skills}


def _node_ingest(state: ClassroomState) -> dict[str, Any]:
    errs = list(state.get("errors") or [])
    out: dict[str, Any] = {"ingest_ok": False}
    slides_bytes = state.get("slides_bytes") or b""
    excel_bytes = state.get("excel_bytes") or b""
    if not slides_bytes or not excel_bytes:
        return {**out, "errors": errs + ["Both slides_bytes and excel_bytes are required."]}
    slides_name = state.get("slides_filename") or "slides.pptx"
    excel_name = state.get("excel_filename") or "responses.xlsx"
    try:
        assert_slides_size(len(slides_bytes))
        assert_excel_size(len(excel_bytes))
        lecture_text, poll_text = parse_slides_bytes(slides_name, slides_bytes)
        df, q_cols = parse_responses_bytes(
            excel_bytes,
            filename=excel_name,
            answer_key=state.get("answer_key"),
        )
        bundle = run_analytics(df, q_cols)
        anonymize = bool(state.get("anonymize", False))
        fig_top = chart_top5(bundle.top5_df, anonymize=anonymize)
        fig_eng = chart_engagement(bundle.responses_df, bundle.question_columns)
        ranked = bundle.ranked_df
        preview_cols = [c for c in ["Student Name", "score_pct", "tier", "rank"] if c in ranked.columns]
        ranked_preview = ranked[preview_cols].head(20).to_dict(orient="records")
        out.update(
            {
                "ingest_ok": True,
                "lecture_text": lecture_text,
                "poll_questions_text": poll_text,
                "tier_counts": bundle.tier_counts,
                "ranked_preview": ranked_preview,
                "question_columns": bundle.question_columns,
                "charts": {
                    "top5": fig_top.to_json(),
                    "engagement": fig_eng.to_json(),
                },
            }
        )
        return out
    except Exception as e:
        return {**out, "errors": errs + [str(e)]}


def _route_after_ingest(state: ClassroomState) -> str:
    if state.get("ingest_ok"):
        return "llm_summary"
    return END


def _extra_system(state: ClassroomState) -> str:
    return combine_agent_skills(
        state.get("agent_context") or "",
        state.get("skills_context") or "",
    )


def _node_summary(state: ClassroomState) -> dict[str, Any]:
    errs = list(state.get("errors") or [])
    if not state.get("want_summary", True):
        return {}
    if not state.get("ingest_ok"):
        return {}
    lecture = (state.get("lecture_text") or "").strip()
    if not lecture:
        return {"errors": errs + ["No text could be extracted from slides for summary."]}
    model = state.get("ollama_model") or "llama3.2"
    client = OllamaClient(model=model)
    extra = _extra_system(state)
    try:
        summary = client.generate_topic_summary(
            lecture,
            state.get("poll_questions_text") or "",
            extra_system=extra or None,
        )
        doc_bytes = build_summary_docx(summary)
        return {"summary_text": summary, "summary_docx_bytes": doc_bytes}
    except Exception as e:
        return {"errors": errs + [f"Summary generation failed: {e}"]}


def _node_homework(state: ClassroomState) -> dict[str, Any]:
    errs = list(state.get("errors") or [])
    if not state.get("want_homework", True):
        return {}
    if not state.get("ingest_ok"):
        return {}
    model = state.get("ollama_model") or "llama3.2"
    client = OllamaClient(model=model)
    topic = (state.get("summary_text") or state.get("lecture_text") or "")[:6000]
    tier_counts = state.get("tier_counts") or {}
    levels = state.get("homework_levels") or ["Extension", "Core", "Support"]
    specs = state.get("question_specs") or [
        {"type": "MCQ", "count": 3},
        {"type": "Fill in the blanks", "count": 2},
        {"type": "Subjective questions", "count": 1},
    ]
    extra = _extra_system(state)
    try:
        homework = client.generate_differentiated_homework(
            topic,
            tier_counts,
            question_specs=specs,
            levels=levels,
            extra_system=extra or None,
        )
        doc_bytes = build_homework_docx(homework)
        return {"homework_text": homework, "homework_docx_bytes": doc_bytes}
    except Exception as e:
        return {"errors": errs + [f"Homework generation failed: {e}"]}


def build_graph():
    g = StateGraph(ClassroomState)
    g.add_node("load_context", _node_load_context)
    g.add_node("ingest", _node_ingest)
    g.add_node("llm_summary", _node_summary)
    g.add_node("llm_homework", _node_homework)

    g.set_entry_point("load_context")
    g.add_edge("load_context", "ingest")
    g.add_conditional_edges("ingest", _route_after_ingest, {"llm_summary": "llm_summary", END: END})
    g.add_edge("llm_summary", "llm_homework")
    g.add_edge("llm_homework", END)

    return g.compile()


def invoke_classroom(
    *,
    slides_bytes: bytes,
    excel_bytes: bytes,
    slides_filename: str = "slides.pptx",
    excel_filename: str = "responses.xlsx",
    answer_key: str | None = None,
    ollama_model: str = "llama3.2",
    want_summary: bool = True,
    want_homework: bool = True,
    anonymize: bool = False,
    homework_levels: list[str] | None = None,
    question_specs: list[dict] | None = None,
) -> ClassroomState:
    """Run the compiled graph and return final state."""
    graph = build_graph()
    init: ClassroomState = {
        "errors": [],
        "slides_bytes": slides_bytes,
        "excel_bytes": excel_bytes,
        "slides_filename": slides_filename,
        "excel_filename": excel_filename,
        "answer_key": answer_key,
        "ollama_model": ollama_model,
        "want_summary": want_summary,
        "want_homework": want_homework,
        "anonymize": anonymize,
        "homework_levels": homework_levels or ["Extension", "Core", "Support"],
        "question_specs": question_specs
        or [
            {"type": "MCQ", "count": 3},
            {"type": "Fill in the blanks", "count": 2},
            {"type": "Subjective questions", "count": 1},
        ],
    }
    result = graph.invoke(init)
    return result  # type: ignore[return-value]

# === api (FastAPI) ===
"""FastAPI: run the LangGraph classroom pipeline."""

import base64
import json
from typing import Any, Optional

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from pydantic import BaseModel, Field


api_app = FastAPI(title="Classroom Report Agent", version="0.2.0")


class HealthResponse(BaseModel):
    status: str
    ollama_ok: bool
    ollama_message: str


class RunResponse(BaseModel):
    ingest_ok: bool
    errors: list[str] = Field(default_factory=list)
    tier_counts: dict[str, int] = Field(default_factory=dict)
    ranked_preview: list[dict[str, Any]] = Field(default_factory=list)
    question_columns: list[str] = Field(default_factory=list)
    charts: Optional[dict[str, str]] = None
    summary_text: Optional[str] = None
    homework_text: Optional[str] = None
    summary_docx_base64: Optional[str] = None
    homework_docx_base64: Optional[str] = None


@api_app.get("/health", response_model=HealthResponse)
def health() -> HealthResponse:
    ok, msg = check_ollama_available()
    return HealthResponse(status="ok", ollama_ok=ok, ollama_message=msg)


@api_app.post("/graph/run", response_model=RunResponse)
@api_app.post("/run", response_model=RunResponse)
@api_app.post("/graph/invoke", response_model=RunResponse)
async def graph_run(
    slides: UploadFile = File(..., description="Lecture slides .pptx or .pdf"),
    responses: UploadFile = File(..., description="Poll responses .xlsx"),
    answer_key: Optional[str] = Form(None),
    ollama_model: str = Form("llama3.2"),
    want_summary: bool = Form(True),
    want_homework: bool = Form(True),
    anonymize: bool = Form(False),
    homework_levels_json: str = Form('["Extension", "Core", "Support"]'),
    question_specs_json: str = Form(
        '[{"type": "MCQ", "count": 3}, {"type": "Fill in the blanks", "count": 2}, {"type": "Subjective questions", "count": 1}]'
    ),
) -> RunResponse:
    try:
        levels = json.loads(homework_levels_json)
        specs = json.loads(question_specs_json)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=400, detail=f"Invalid JSON in form fields: {e}") from e

    s_bytes = await slides.read()
    e_bytes = await responses.read()
    if not s_bytes or not e_bytes:
        raise HTTPException(status_code=400, detail="Slides and responses files must not be empty.")

    state = invoke_classroom(
        slides_bytes=s_bytes,
        excel_bytes=e_bytes,
        slides_filename=slides.filename or "slides.pptx",
        excel_filename=responses.filename or "responses.xlsx",
        answer_key=answer_key,
        ollama_model=ollama_model,
        want_summary=want_summary,
        want_homework=want_homework,
        anonymize=anonymize,
        homework_levels=levels if isinstance(levels, list) else None,
        question_specs=specs if isinstance(specs, list) else None,
    )

    sum_b64 = None
    hw_b64 = None
    if state.get("summary_docx_bytes"):
        sum_b64 = base64.standard_b64encode(state["summary_docx_bytes"]).decode("ascii")
    if state.get("homework_docx_bytes"):
        hw_b64 = base64.standard_b64encode(state["homework_docx_bytes"]).decode("ascii")

    return RunResponse(
        ingest_ok=bool(state.get("ingest_ok")),
        errors=list(state.get("errors") or []),
        tier_counts=dict(state.get("tier_counts") or {}),
        ranked_preview=list(state.get("ranked_preview") or []),
        question_columns=list(state.get("question_columns") or []),
        charts=state.get("charts"),
        summary_text=state.get("summary_text"),
        homework_text=state.get("homework_text"),
        summary_docx_base64=sum_b64,
        homework_docx_base64=hw_b64,
    )


# === CLI ===
def cli_main() -> None:
    p = argparse.ArgumentParser(description="Classroom report agent (LangGraph + Ollama)")
    p.add_argument("slides", type=Path, help="Path to lecture slides (.pptx or .pdf)")
    p.add_argument("excel", type=Path, help="Path to poll responses (.xlsx)")
    p.add_argument("--answer-key", default=None, help="Comma-separated correct answers")
    p.add_argument("--ollama-model", default="llama3.2")
    p.add_argument("--no-summary", action="store_true")
    p.add_argument("--no-homework", action="store_true")
    p.add_argument("--anonymize", action="store_true")
    p.add_argument("--out-dir", type=Path, default=Path("."), help="Where to write .docx outputs")
    p.add_argument("--homework-levels", default='["Extension", "Core", "Support"]', help="JSON array of levels")
    p.add_argument(
        "--question-specs",
        default='[{"type": "MCQ", "count": 3}, {"type": "Fill in the blanks", "count": 2}, {"type": "Subjective questions", "count": 1}]',
        help="JSON array of {type, count} objects",
    )
    args = p.parse_args()
    if not args.slides.is_file():
        print(f"Slides not found: {args.slides}", file=sys.stderr)
        sys.exit(1)
    if not args.excel.is_file():
        print(f"Excel not found: {args.excel}", file=sys.stderr)
        sys.exit(1)
    try:
        levels = json.loads(args.homework_levels)
        specs = json.loads(args.question_specs)
    except json.JSONDecodeError as e:
        print(f"Invalid JSON: {e}", file=sys.stderr)
        sys.exit(1)
    state = invoke_classroom(
        slides_bytes=args.slides.read_bytes(),
        excel_bytes=args.excel.read_bytes(),
        slides_filename=args.slides.name,
        excel_filename=args.excel.name,
        answer_key=args.answer_key,
        ollama_model=args.ollama_model,
        want_summary=not args.no_summary,
        want_homework=not args.no_homework,
        anonymize=args.anonymize,
        homework_levels=levels if isinstance(levels, list) else None,
        question_specs=specs if isinstance(specs, list) else None,
    )
    out = args.out_dir
    out.mkdir(parents=True, exist_ok=True)
    errs = state.get("errors") or []
    for e in errs:
        print(f"Error: {e}", file=sys.stderr)
    if state.get("summary_docx_bytes"):
        pth = out / "class_topic_summary.docx"
        pth.write_bytes(state["summary_docx_bytes"])
        print(f"Wrote {pth}")
    if state.get("homework_docx_bytes"):
        pth = out / "differentiated_homework.docx"
        pth.write_bytes(state["homework_docx_bytes"])
        print(f"Wrote {pth}")
    if state.get("charts"):
        cd = out / "charts"
        cd.mkdir(parents=True, exist_ok=True)
        for name, fig_json in (state.get("charts") or {}).items():
            (cd / f"{name}.json").write_text(fig_json, encoding="utf-8")
        print(f"Wrote Plotly JSON under {cd}")
    preview = state.get("ranked_preview") or []
    if preview:
        print("\nRanked preview (first rows):", json.dumps(preview[:5], indent=2))
    sys.exit(1 if errs else 0 if state.get("ingest_ok") else 1)



# === Streamlit ===
def run_streamlit() -> None:
    import streamlit as st

    # Page config
    st.set_page_config(
        page_title="Classroom Report & Analytics",
        page_icon="📚",
        layout="wide",
        initial_sidebar_state="expanded",
    )

    # Guardrail disclaimer (always visible)
    st.caption(
        "Analytics and tiers are based only on poll responses; use professional judgment when assigning homework."
    )


    def _parse_slides(slide_file):
        """Return (lecture_text, poll_questions_text) from uploaded slide file."""
        if slide_file is None:
            return "", ""
        name = slide_file.name or "slides.pptx"
        return parse_slides_bytes(name, slide_file.getvalue())


    def _parse_responses(excel_file, answer_key):
        """Load Excel and return (df, question_columns). Raises on error."""
        return parse_responses_bytes(
            excel_file.getvalue(),
            filename=excel_file.name or "responses.xlsx",
            answer_key=answer_key,
        )


    def _ensure_analytics_data():
        """Parse uploads and compute analytics; store in session_state. Returns error message or None."""
        if st.session_state.get("_analytics_ready"):
            return None
        excel_file = _get_excel_file()
        if not excel_file:
            return "Upload an Excel file first."
        try:
            size = len(excel_file.getvalue()) if hasattr(excel_file, "getvalue") else getattr(excel_file, "size", 0)
        except Exception:
            size = 0
        if size > MAX_EXCEL_SIZE_BYTES:
            return "Excel file too large."
        try:
            if hasattr(excel_file, "seek"):
                excel_file.seek(0)
            df, q_cols = _parse_responses(excel_file, st.session_state.get("answer_key"))
        except Exception as e:
            return str(e)
        bundle = run_analytics(df, q_cols)
        scores_df = bundle.scores_df
        ranked_df = bundle.ranked_df
        top5_df = bundle.top5_df
        tier_counts = bundle.tier_counts
        st.session_state["responses_df"] = df
        st.session_state["question_columns"] = q_cols
        st.session_state["scores_df"] = scores_df
        st.session_state["ranked_df"] = ranked_df
        st.session_state["top5_df"] = top5_df
        st.session_state["tier_counts"] = tier_counts
        st.session_state["_analytics_ready"] = True
        return None


    def _ensure_lecture_data():
        """Parse slides and store lecture_text, poll_questions_text. Returns error or None."""
        if st.session_state.get("_lecture_ready"):
            return None
        slide_file = _get_slide_file()
        if not slide_file:
            return "Upload slides first."
        if slide_file.size > MAX_SLIDES_SIZE_BYTES:
            return "Slides file too large."
        try:
            lecture_text, poll_text = _parse_slides(slide_file)
        except Exception as e:
            return str(e)
        st.session_state["lecture_text"] = lecture_text
        st.session_state["poll_questions_text"] = poll_text
        st.session_state["_lecture_ready"] = True
        return None


    # Sidebar: config
    with st.sidebar:
        st.header("Settings")
        ollama_model = st.text_input("Ollama model", value="llama3.2", key="ollama_model", help="Model pulled via ollama pull <name>")
        anonymize = st.checkbox("Anonymize in report", value=False, key="anonymize", help="Use 'Student 1' etc. instead of names in outputs")
        st.divider()
        st.subheader("File limits")
        st.write(f"Slides: max {MAX_SLIDES_SIZE_BYTES // (1024*1024)} MB")
        st.write(f"Excel: max {MAX_EXCEL_SIZE_BYTES // (1024*1024)} MB")

    # Navigation: pages
    st.title("Classroom Report & Analytics")

    page = st.radio(
        "Go to",
        ["Upload", "Analytics", "Reports"],
        horizontal=True,
        label_visibility="collapsed",
        key="page_radio",
    )


    def _get_slide_file():
        """Return the uploaded slide file (from widget or from persisted bytes)."""
        f = st.session_state.get("slide_file")
        if f is not None:
            try:
                f.getvalue()
                return f
            except Exception:
                pass
        raw = st.session_state.get("slide_file_bytes")
        name = st.session_state.get("slide_file_name", "slides.pptx")
        if raw is not None:
            buf = BytesIO(raw)
            buf.name = name
            return buf
        return None


    def _get_excel_file():
        """Return the uploaded Excel file (from widget or from persisted bytes)."""
        f = st.session_state.get("excel_file")
        if f is not None:
            try:
                f.getvalue()
                return f
            except Exception:
                try:
                    f.seek(0)
                    return f
                except Exception:
                    pass
        raw = st.session_state.get("excel_file_bytes")
        name = st.session_state.get("excel_file_name", "responses.xlsx")
        if raw is not None:
            buf = BytesIO(raw)
            buf.name = name
            return buf
        return None


    if page == "Upload":
        st.header("Upload")
        slide_file = st.file_uploader(
            "Lecture slides (PPT or PDF)",
            type=[e.lstrip(".") for e in ALLOWED_SLIDE_EXTENSIONS],
            help="One file per lecture. Poll questions can be on slides titled 'Poll' or 'Question'.",
            key="slide_upload",
        )
        excel_file = st.file_uploader(
            "Poll responses (Excel)",
            type=[e.lstrip(".") for e in ALLOWED_EXCEL_EXTENSIONS],
            help="Student Name + Q1, Q2, ... (1/0) OR Student Name + Q1_Selected/Q1_Correct, Q2_Selected/Q2_Correct, ... (letters A–D).",
            key="excel_upload",
        )
        answer_key = st.text_input(
            "Correct answers (optional)",
            placeholder="e.g. 1,1,0,1 or A,B,A,C",
            help="Comma-separated; one value per question. If empty, we assume 1=correct.",
            key="answer_key_input",
        )
        if slide_file:
            if slide_file.size > MAX_SLIDES_SIZE_BYTES:
                st.error(f"Slides file too large. Max {MAX_SLIDES_SIZE_BYTES // (1024*1024)} MB.")
            else:
                st.success(f"Slides: {slide_file.name} ({slide_file.size // 1024} KB)")
        if excel_file:
            if excel_file.size > MAX_EXCEL_SIZE_BYTES:
                st.error(f"Excel file too large. Max {MAX_EXCEL_SIZE_BYTES // (1024*1024)} MB.")
            else:
                st.success(f"Responses: {excel_file.name} ({excel_file.size // 1024} KB)")

        # Persist in session state for other pages; save bytes so uploads survive when user switches to Analytics/Reports
        st.session_state["slide_file"] = slide_file
        st.session_state["excel_file"] = excel_file
        if slide_file is not None:
            st.session_state["slide_file_bytes"] = slide_file.getvalue()
            st.session_state["slide_file_name"] = slide_file.name
        if excel_file is not None:
            st.session_state["excel_file_bytes"] = excel_file.getvalue()
            st.session_state["excel_file_name"] = excel_file.name
        st.session_state["answer_key"] = answer_key.strip() if answer_key else None
        # Invalidate caches when files change
        fingerprint = (
            (slide_file.name, slide_file.size) if slide_file else (None, None),
            (excel_file.name, excel_file.size) if excel_file else (None, None),
        )
        if st.session_state.get("_upload_fingerprint") != fingerprint:
            st.session_state["_upload_fingerprint"] = fingerprint
            st.session_state.pop("_lecture_ready", None)
            st.session_state.pop("_analytics_ready", None)

    elif page == "Analytics":
        st.header("Analytics")
        err = _ensure_analytics_data()
        if err:
            st.error(err)
            st.info("Upload an Excel file with poll responses on the Upload page first.")
        else:
            responses_df = st.session_state["responses_df"]
            question_columns = st.session_state["question_columns"]
            top5_df = st.session_state["top5_df"]
            ranked_df = st.session_state["ranked_df"]
            tier_counts = st.session_state["tier_counts"]

            c1, c2 = st.columns(2)
            with c1:
                fig_top5 = chart_top5(top5_df, anonymize=st.session_state.get("anonymize", False))
                st.plotly_chart(fig_top5, use_container_width=True)
            with c2:
                fig_eng = chart_engagement(responses_df, question_columns)
                st.plotly_chart(fig_eng, use_container_width=True)

            st.subheader("Tier counts")
            st.write("Extension (top): ", tier_counts.get("top", 0), " | Core (average): ", tier_counts.get("average", 0), " | Support (low): ", tier_counts.get("low", 0))
            st.dataframe(ranked_df[["Student Name", "score_pct", "tier", "rank"]].head(20), use_container_width=True, hide_index=True)

    elif page == "Reports":
        st.header("Reports")
        if _get_slide_file() is None or _get_excel_file() is None:
            st.info("Upload both slides and poll responses on the Upload page first.")
        else:
            client = OllamaClient(model=st.session_state.get("ollama_model", "llama3.2"))
            ok, msg = client.available()
            if not ok:
                st.warning(msg)
            else:
                st.success("Ollama is available.")

            _ensure_lecture_data()
            _ensure_analytics_data()
            lecture_text = st.session_state.get("lecture_text", "")
            poll_questions_text = st.session_state.get("poll_questions_text", "")
            tier_counts = st.session_state.get("tier_counts", {})

            # Generate Summary
            st.subheader("Topic summary (Word)")
            if st.button("Generate summary report"):
                if not ok:
                    st.error("Start Ollama and pull a model first.")
                elif not lecture_text:
                    st.error("No text could be extracted from the slides.")
                else:
                    with st.spinner("Generating summary..."):
                        try:
                            summary = client.generate_topic_summary(lecture_text, poll_questions_text)
                            doc_bytes = build_summary_docx(summary)
                            st.session_state["summary_docx_bytes"] = doc_bytes
                            st.session_state["summary_text"] = summary
                            st.session_state["summary_generated"] = True
                            st.success("Summary generated.")
                        except Exception as e:
                            st.error(str(e))

            if st.session_state.get("summary_generated") and st.session_state.get("summary_docx_bytes"):
                st.download_button(
                    "Download summary report (.docx)",
                    data=st.session_state["summary_docx_bytes"],
                    file_name="class_topic_summary.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    key="dl_summary",
                )

            st.divider()
            st.subheader("Differentiated homework (Word)")
            if not st.session_state.get("summary_text") and st.session_state.get("lecture_text"):
                st.caption("Tip: Generate the topic summary first so homework activities are based on your lecture content.")

            # Level(s) to generate
            st.markdown("**Generate for levels:**")
            levels_options = ["Extension", "Core", "Support"]
            selected_levels = st.multiselect(
                "Choose which level(s) to generate questions for",
                options=levels_options,
                default=levels_options,
                key="hw_levels",
            )
            if not selected_levels:
                st.caption("Select at least one level (Extension, Core, and/or Support).")

            # Homework options: question types and counts per type
            st.markdown("**Activity options** — question types and how many per selected level:")
            col1, col2, col3 = st.columns(3)
            question_specs = []
            with col1:
                include_mcq = st.checkbox("MCQ", value=True, key="hw_include_mcq")
                num_mcq = st.number_input("Number", min_value=1, max_value=15, value=3, key="hw_num_mcq") if include_mcq else 0
                if include_mcq:
                    question_specs.append({"type": "MCQ", "count": num_mcq})
            with col2:
                include_fib = st.checkbox("Fill in the blanks", value=True, key="hw_include_fib")
                num_fib = st.number_input("Number", min_value=1, max_value=15, value=2, key="hw_num_fib") if include_fib else 0
                if include_fib:
                    question_specs.append({"type": "Fill in the blanks", "count": num_fib})
            with col3:
                include_subj = st.checkbox("Subjective questions", value=True, key="hw_include_subj")
                num_subj = st.number_input("Number", min_value=1, max_value=15, value=1, key="hw_num_subj") if include_subj else 0
                if include_subj:
                    question_specs.append({"type": "Subjective questions", "count": num_subj})
            if not question_specs:
                st.caption("Select at least one question type above.")

            if st.button("Generate homework"):
                if not ok:
                    st.error("Start Ollama and pull a model first.")
                elif not selected_levels:
                    st.error("Select at least one level (Extension, Core, and/or Support).")
                elif not question_specs:
                    st.error("Select at least one question type and set the number of questions.")
                else:
                    with st.spinner("Generating homework..."):
                        try:
                            topic = st.session_state.get("summary_text") or st.session_state.get("lecture_text", "")
                            topic = (topic or "")[:6000]
                            homework = client.generate_differentiated_homework(
                                topic, tier_counts, question_specs=question_specs, levels=selected_levels
                            )
                            doc_bytes = build_homework_docx(homework)
                            st.session_state["homework_docx_bytes"] = doc_bytes
                            st.session_state["homework_generated"] = True
                            st.success("Homework generated.")
                        except Exception as e:
                            st.error(str(e))

            if st.session_state.get("homework_generated") and st.session_state.get("homework_docx_bytes"):
                st.download_button(
                    "Download homework (.docx)",
                    data=st.session_state["homework_docx_bytes"],
                    file_name="differentiated_homework.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    key="dl_homework",
                )

if __name__ == "__main__":
    if len(sys.argv) > 1 and sys.argv[1] == "cli":
        sys.argv = [sys.argv[0]] + sys.argv[2:]
        cli_main()
    else:
        run_streamlit()
