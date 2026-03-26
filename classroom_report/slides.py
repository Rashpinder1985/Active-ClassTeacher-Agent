"""Extract text from .pptx and .pdf; detect poll-like slides."""

from __future__ import annotations

import tempfile
from io import BytesIO
from pathlib import Path
from typing import List, Tuple, Union

from pptx import Presentation

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None  # type: ignore


def extract_text_from_pptx(path: Union[str, Path, BytesIO]) -> Tuple[str, List[Tuple[int, str]]]:
    if isinstance(path, BytesIO):
        path = BytesIO(path.getvalue())
        path.name = "upload.pptx"
    p = Path(path) if not isinstance(path, BytesIO) else None
    if p and p.suffix.lower() != ".pptx":
        raise ValueError("File must be .pptx")
    if isinstance(path, BytesIO):
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(path.getvalue())
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
        finally:
            Path(tmp_path).unlink(missing_ok=True)
    else:
        prs = Presentation(str(path))
    slides_text: List[Tuple[int, str]] = []
    all_parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        slide_text = "\n".join(parts)
        slides_text.append((i, slide_text))
        if slide_text:
            all_parts.append(slide_text)
    return "\n\n".join(all_parts), slides_text


def extract_text_from_pdf(path: Union[str, Path, BytesIO]) -> Tuple[str, List[Tuple[int, str]]]:
    if fitz is None:
        raise ImportError("PyMuPDF is required for PDF support. Install with: pip install PyMuPDF")
    if isinstance(path, BytesIO):
        doc = fitz.open(stream=path.getvalue(), filetype="pdf")
    else:
        path = Path(path)
        if path.suffix.lower() != ".pdf":
            raise ValueError("File must be .pdf")
        doc = fitz.open(str(path))
    slides_text: List[Tuple[int, str]] = []
    all_parts: List[str] = []
    try:
        for i in range(len(doc)):
            text = doc[i].get_text().strip()
            slides_text.append((i + 1, text))
            if text:
                all_parts.append(text)
    finally:
        doc.close()
    return "\n\n".join(all_parts), slides_text


def get_poll_slides(slides_text: List[Tuple[int, str]]) -> List[Tuple[int, str]]:
    poll_keywords = ("poll", "question", "quiz")
    result = []
    for idx, text in slides_text:
        first_line = (text.split("\n")[0] or "").strip().lower()
        if any(kw in first_line for kw in poll_keywords):
            result.append((idx, text))
    return result
